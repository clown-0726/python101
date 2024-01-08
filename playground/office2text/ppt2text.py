from pptx import Presentation

'''
pip install python-pptx

参考连接：https://zhuanlan.zhihu.com/p/405916351
'''

pptx = Presentation('./asset/orbit_strategy.pptx')

for slide in pptx.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                print(paragraph.text)

if __name__ == '__main__':
    pass
